from pptx import Presentation

from services.layout_template_service import get_template_for_slide
import services.ppt_render_service as render_service
from services.ppt_style_config import get_ppt_style_profile


def _task():
    return {
        "id": 20011,
        "course_title": "[TEST] Round20.1 Template Variant",
        "topic": "The Internet",
        "grade": "高一",
        "textbook": "人教版",
        "volume": "必修二",
        "unit": "Unit 3",
        "lesson_type": "Reading",
        "duration": 45,
        "student_level": "中等",
        "style": "常规课",
    }


def _new_slide():
    prs = Presentation()
    prs.slide_width = render_service.Inches(render_service.SLIDE_SIZE[0])
    prs.slide_height = render_service.Inches(render_service.SLIDE_SIZE[1])
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def run():
    mapping = {
        "objectives": "numbered_cards",
        "lead_in": "question_cards",
        "prediction": "three_step_flow",
        "careful_reading": "question_evidence_columns",
        "vocabulary_focus": "vocabulary_cards",
        "summary": "three_takeaway_cards",
        "homework": "homework_levels",
    }
    for slide_type, variant in mapping.items():
        tpl = get_template_for_slide(slide_type)
        assert tpl.get("visual_variant") == variant

    calls = []

    def _make_stub(name):
        def _stub(slide, slide_data, task, prs, style, template):
            calls.append(name)
        return _stub

    backup = dict(render_service.TEMPLATE_VARIANT_RENDERERS)
    try:
        for name in set(mapping.values()):
            render_service.TEMPLATE_VARIANT_RENDERERS[name] = _make_stub(name)
        style = get_ppt_style_profile(_task())
        for slide_type, variant in mapping.items():
            prs, slide = _new_slide()
            slide_data = {
                "slide_type": slide_type,
                "title": slide_type,
                "visible_content": ["a", "b", "c"],
            }
            tpl = get_template_for_slide(slide_type)
            result = render_service.render_by_layout_template(slide, slide_data, _task(), tpl, prs, style)
            assert result == variant
            assert calls[-1] == variant
            assert result != "plain_fallback"
    finally:
        render_service.TEMPLATE_VARIANT_RENDERERS.clear()
        render_service.TEMPLATE_VARIANT_RENDERERS.update(backup)

    # old style compatibility path: no visual variant falls back to plain renderer
    prs, slide = _new_slide()
    style = get_ppt_style_profile(_task())
    fallback_template = {
        "template_name": "fallback_template",
        "max_blocks": 2,
        "max_items_per_block": 2,
        "body_font_size": 18,
        "min_font_size": 14,
    }
    result = render_service.render_by_layout_template(
        slide,
        {"slide_type": "other", "title": "Other", "visible_content": ["x", "y", "z"]},
        _task(),
        fallback_template,
        prs,
        style,
    )
    assert result == "plain_fallback"

    print("ROUND_0201_OK")


if __name__ == "__main__":
    run()
