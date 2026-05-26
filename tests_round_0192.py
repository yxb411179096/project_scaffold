from services.ppt_graphics_service import (
    GRAPHIC_MIN_FONT,
    render_mindmap,
    render_task_steps,
    render_vocabulary_cards,
    smart_clip,
)
from services.ppt_quality_check_service import check_ppt_quality
from services.ppt_render_service import export_pptx
from pptx import Presentation


def _task():
    return {
        "id": 19021,
        "course_title": "[TEST] Round192 Clip",
        "topic": "The Internet",
        "grade": "高一",
        "textbook": "人教版",
        "volume": "必修二",
        "unit": "Unit 3",
        "lesson_type": "Reading",
        "duration": 45,
        "student_level": "中等",
        "style": "常规课",
    }


def run():
    # 1 smart_clip should not hard-cut English half-word
    text = "Review the key language and the main classroom task."
    clipped = smart_clip(text, 38, preserve_words=True)
    assert clipped.endswith("...")
    assert "classr..." not in clipped

    # 2/3/4 density limits for cards/branches/steps
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = {
        "primary": (21, 76, 121),
        "secondary": (51, 115, 167),
        "accent": (244, 178, 102),
        "text": (32, 43, 57),
        "muted": (102, 112, 133),
        "surface": (255, 255, 255),
        "soft": (240, 246, 252),
    }
    render_vocabulary_cards(
        slide,
        prs,
        {"cards": [{"word": f"Very long word item content {i} should be converted", "meaning": "meaning" * 10, "example": "example " * 20} for i in range(6)]},
        area=(0.8, 1.2, 10.8, 4.8),
        style=style,
    )
    render_mindmap(
        slide,
        prs,
        {"center": "Summary", "branches": [f"branch {i} with too much text for a node" for i in range(7)]},
        area=(0.8, 1.2, 10.8, 4.8),
        style=style,
    )
    render_task_steps(
        slide,
        prs,
        {"steps": [f"step {i} with long instruction text " * 4 for i in range(6)]},
        area=(0.8, 1.2, 10.8, 4.8),
        style=style,
    )

    # 5 min font constant safeguard
    assert GRAPHIC_MIN_FONT >= 9

    # 6 quality check graphic_too_dense
    report = check_ppt_quality(
        _task(),
        [
            {
                "slide_index": 1,
                "slide_type": "careful_reading",
                "title": "Careful Reading",
                "visible_content": ["line 1", "line 2", "line 3", "line 4", "line 5"],
                "teacher_notes": "This teacher note is long enough to pass minimum length check.",
                "graphic_type": "reading_structure",
                "graphic_data": {
                    "nodes": [{"label": f"L{i}", "content": "x" * 100} for i in range(9)]
                },
                "layout_plan": {"layout_type": "reading_task_layout", "graphic_type": "reading_structure"},
            }
        ],
    )
    codes = {
        d["code"]
        for s in report["slides"]
        for d in s.get("issue_details", [])
    }
    assert "graphic_too_dense" in codes

    # 7 old task still export
    path = export_pptx(
        _task(),
        [
            {
                "slide_index": 1,
                "slide_type": "cover",
                "title": "Legacy Slide",
                "visible_content": ["Unit 3", "Reading"],
                "teacher_notes": "notes",
                "teaching_purpose": "open",
                "estimated_time": "2 min",
                "interaction_type": "Teacher-student interaction",
                "layout_plan": {"layout_type": "cover_layout"},
            }
        ],
    )
    assert path.exists()

    print("ROUND_0192_OK")


if __name__ == "__main__":
    run()
