from app import create_app

from services.layout_template_service import get_template_for_slide
from services.ppt_render_service import render_by_layout_template
from services.ppt_style_config import get_ppt_style_profile
from services.ppt_style_service import (
    get_style_options,
    get_style_preset,
    recommend_style_key,
)
from pptx import Presentation


def _task(lesson_type="Reading", ppt_style=None):
    task = {
        "id": 21001,
        "course_title": "[TEST] Round21 Style",
        "topic": "The Internet",
        "grade": "高一",
        "textbook": "人教版",
        "volume": "必修二",
        "unit": "Unit 3",
        "lesson_type": lesson_type,
        "duration": 45,
        "student_level": "中等",
        "style": "常规课",
    }
    if ppt_style:
        task["ppt_style"] = ppt_style
    return task


def run():
    # 1 / 2 style service base lookups
    default = get_style_preset("default")
    open_class = get_style_preset("open_class")
    assert default.get("style_key") == "default"
    assert open_class.get("style_key") == "open_class"

    # 3 / 4 / 5 recommendation
    assert recommend_style_key("Reading", "常规课") == "reading_focus"
    assert recommend_style_key("Writing", "常规课") == "writing_focus"
    assert recommend_style_key("Grammar", "常规课") == "grammar_focus"

    # 6 old task fallback style
    legacy = get_ppt_style_profile(_task("Reading"))
    assert legacy.get("name") == "reading_focus"

    # 7 render layer reads style preset
    profile = get_ppt_style_profile(_task("Reading", ppt_style="open_class"))
    assert profile.get("name") == "open_class"
    assert profile.get("display_name")

    # 8 visual variant renderer receives style profile
    from services import ppt_render_service as render_service

    calls = []

    def _stub(slide, slide_data, task, prs, style, template):
        calls.append(style.get("name"))

    backup = dict(render_service.TEMPLATE_VARIANT_RENDERERS)
    try:
        render_service.TEMPLATE_VARIANT_RENDERERS["numbered_cards"] = _stub
        prs = Presentation()
        prs.slide_width = render_service.Inches(render_service.SLIDE_SIZE[0])
        prs.slide_height = render_service.Inches(render_service.SLIDE_SIZE[1])
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tpl = get_template_for_slide("objectives")
        render_by_layout_template(
            slide,
            {"slide_type": "objectives", "title": "Objectives", "visible_content": ["A", "B", "C"]},
            _task("Reading", ppt_style="open_class"),
            tpl,
            prs,
            get_ppt_style_profile(_task("Reading", ppt_style="open_class")),
        )
    finally:
        render_service.TEMPLATE_VARIANT_RENDERERS.clear()
        render_service.TEMPLATE_VARIANT_RENDERERS.update(backup)
    assert calls and calls[-1] == "open_class"

    # 9 / 10 pages expose ppt_style field
    app = create_app()
    app.config["TESTING"] = True
    with app.test_client() as client:
        new_html = client.get("/ppt/new").get_data(as_text=True)
        manuscript_html = client.get("/ppt/from-manuscript").get_data(as_text=True)
    assert 'name="ppt_style"' in new_html
    assert 'name="ppt_style"' in manuscript_html

    # ensure option set includes new style family
    keys = {k for k, _ in get_style_options()}
    assert {"default", "open_class", "fresh_classroom", "reading_focus", "writing_focus", "grammar_focus"} <= keys

    print("ROUND_021_OK")


if __name__ == "__main__":
    run()
