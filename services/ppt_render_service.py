"""PPTX renderer using only python-pptx public APIs.

The current exporter intentionally avoids speaker notes and any direct OOXML
manipulation so the generated deck stays compatible with PowerPoint, Keynote,
and WPS.
"""

import zipfile
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from config import EXPORT_PPTX_DIR
from services.agents.layout_planner_agent import plan_layout_for_slide
from services.export_utils import safe_course_filename
from services.layout_template_service import get_template_for_slide
from services.ppt_graphics_service import GRAPHIC_RENDERERS
from services.ppt_style_config import get_ppt_style_profile


SLIDE_SIZE = (13.333, 7.5)

# Unified design constants (0.18.0)
SLIDE_W = SLIDE_SIZE[0]
SLIDE_H = SLIDE_SIZE[1]
MARGIN_X = 0.9
MARGIN_Y = 0.55
SAFE_BOTTOM = 6.85
FOOTER_HEIGHT = 0.28
TITLE_FONT_SIZE = 24
SUBTITLE_FONT_SIZE = 16
BODY_FONT_SIZE = 16
SMALL_FONT_SIZE = 10
CARD_RADIUS = 0.12
CARD_GAP = 0.2
MAX_BULLETS = 5
MAX_BULLET_CHARS = 80
MAX_TITLE_CHARS = 90

COLOR_PRIMARY = (21, 76, 121)
COLOR_SECONDARY = (51, 115, 167)
COLOR_ACCENT = (244, 178, 102)
COLOR_TEXT = (32, 43, 57)
COLOR_MUTED = (102, 112, 133)
COLOR_CARD_BG = (248, 251, 253)
COLOR_BORDER = (215, 225, 234)
COLOR_WARNING = (215, 94, 46)
GRAPHIC_AREA = (0.95, 1.85, 11.2, 4.25)
MAX_MAIN_BLOCKS = 3


def _rgb(color_tuple):
    return RGBColor(*color_tuple)


def _set_run_style(paragraph, size, color, bold=False, italic=False):
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = _rgb(color)
    paragraph.font.bold = bold
    paragraph.font.italic = italic


def estimate_text_lines(text, max_chars_per_line):
    text = str(text or "").strip()
    if not text:
        return 1
    max_chars_per_line = max(8, int(max_chars_per_line or 30))
    lines = 0
    for chunk in text.splitlines() or [text]:
        chunk = chunk.strip()
        if not chunk:
            lines += 1
            continue
        lines += max(1, (len(chunk) + max_chars_per_line - 1) // max_chars_per_line)
    return lines


def truncate_text(text, max_chars, suffix="..."):
    text = str(text or "").strip()
    if max_chars <= 0:
        return ""
    if len(text) <= max_chars:
        return text
    return text[: max(1, max_chars - len(suffix))].rstrip() + suffix


def fit_font_size(text, box_width, box_height, max_size, min_size):
    width = max(float(box_width), 0.2)
    height = max(float(box_height), 0.2)
    lines = estimate_text_lines(text, max_chars_per_line=max(12, int(width * 11)))
    size = int(min(max_size, (height * 72) / max(lines, 1)))
    return max(min_size, size)


def split_long_text_to_bullets(text):
    text = str(text or "").strip()
    if not text:
        return []
    parts = re.split(r"[。！？!?;；\n]+", text)
    return [p.strip() for p in parts if p and p.strip()]


def normalize_bullets(items, max_items=MAX_BULLETS, max_chars_each=MAX_BULLET_CHARS):
    normalized = []
    for item in list(items or []):
        text = str(item or "").strip()
        if not text:
            continue
        if len(text) > max_chars_each:
            text = truncate_text(text, max_chars_each)
        normalized.append(text)
        if len(normalized) >= max_items:
            break
    return normalized


def _set_background(slide, prs, style):
    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        prs.slide_height,
    )
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb(style["background"])
    background.line.fill.background()


def _add_panel(slide, left, top, width, height, fill_color, line_color=None, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    panel = slide.shapes.add_shape(shape_type, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(fill_color)
    if line_color:
        panel.line.color.rgb = _rgb(line_color)
    else:
        panel.line.fill.background()
    return panel


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    size=18,
    color=None,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = valign
    text_frame.text = truncate_text(text, 420)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    _set_run_style(paragraph, size, color or (32, 43, 57), bold=bold)
    return textbox


def _add_bullets(slide, items, left, top, width, height, style, size=20, level=0):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    for index, item in enumerate(normalize_bullets(items, max_items=MAX_BULLETS, max_chars_each=MAX_BULLET_CHARS)):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = str(item)
        paragraph.level = level
        paragraph.space_after = Pt(8)
        _set_run_style(paragraph, size, style["text"])
    return textbox


def _add_header_band(slide, slide_data, style):
    _add_panel(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.85), style["primary"])
    _add_textbox(
        slide,
        Inches(0.95),
        Inches(0.56),
        Inches(8.9),
        Inches(0.35),
        text=slide_data["title"],
        size=TITLE_FONT_SIZE,
        color=(255, 255, 255),
        bold=True,
    )


def _add_footer(slide, slide_index, task, style):
    footer_text = f"{task.get('unit') or task.get('course_title') or 'Lesson'}"
    _add_textbox(
        slide,
        Inches(MARGIN_X),
        Inches(SAFE_BOTTOM),
        Inches(8.0),
        Inches(FOOTER_HEIGHT),
        text=truncate_text(footer_text, 72),
        size=SMALL_FONT_SIZE,
        color=style.get("muted", COLOR_MUTED),
    )
    _add_panel(slide, Inches(11.85), Inches(6.78), Inches(0.78), Inches(0.33), style["soft"])
    _add_textbox(
        slide,
        Inches(11.96),
        Inches(6.83),
        Inches(0.56),
        Inches(0.16),
        text=str(slide_index),
        size=SMALL_FONT_SIZE,
        color=style["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _limit_items(items, limit=5):
    cleaned = [truncate_text(str(item).strip(), MAX_BULLET_CHARS) for item in items if str(item or "").strip()]
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 1] + ["More details are included in the teacher script."]


def _is_orphan_number_marker(text):
    return bool(re.match(r"^(?:\d+|[A-Za-z])[.)、．]$", str(text or "").strip()))


def _get_layout_plan(slide_data):
    return slide_data.get("layout_plan") or {}


def _get_content_blocks(slide_data):
    layout_plan = _get_layout_plan(slide_data)
    blocks = layout_plan.get("content_blocks") or []
    return blocks


def _find_block(blocks, role):
    for block in blocks:
        if block.get("role") == role:
            return block
    return None


def _block_title(block, default):
    return (block or {}).get("title") or default


def _block_items(block, fallback=None):
    source_items = (block or {}).get("items") or fallback or []
    if isinstance(source_items, str):
        source_items = split_long_text_to_bullets(source_items)
    items = normalize_bullets(source_items, max_items=MAX_BULLETS, max_chars_each=MAX_BULLET_CHARS)
    banned = {
        "complete the classroom task.",
        "complete the task.",
        "please complete the task without specific instruction.",
    }
    cleaned = []
    for item in items:
        text = str(item or "").strip()
        lowered = text.lower()
        if lowered in banned:
            continue
        if re.fullmatch(r"step\s*[1-3]\s*[:：]?\s*$", lowered):
            continue
        cleaned.append(text)
    return cleaned or ["Follow the task and prepare one clear answer with evidence."]


def _add_teacher_hint(slide, slide_data, style):
    layout_plan = _get_layout_plan(slide_data)
    position = layout_plan.get("teacher_hint_position", "footer")
    text = f"Purpose: {slide_data.get('teaching_purpose', '')}"

    if position == "right_panel":
        _add_textbox(
            slide,
            Inches(9.55),
            Inches(5.55),
            Inches(2.1),
            Inches(0.5),
            text=text,
            size=9,
            color=style["muted"],
        )
    elif position == "bottom_right":
        _add_textbox(
            slide,
            Inches(7.9),
            Inches(6.2),
            Inches(3.6),
            Inches(0.32),
            text=text,
            size=9,
            color=style["muted"],
        )
    elif position == "bottom_left":
        _add_textbox(
            slide,
            Inches(0.9),
            Inches(6.2),
            Inches(4.4),
            Inches(0.32),
            text=text,
            size=9,
            color=style["muted"],
        )
    else:
        _add_textbox(
            slide,
            Inches(0.95),
            Inches(6.18),
            Inches(10.2),
            Inches(0.3),
            text=text,
            size=9,
            color=style["muted"],
        )


def _resolve_graphic_type(slide_data):
    direct = str(slide_data.get("graphic_type") or "").strip()
    if direct:
        return direct
    layout_graphic = str((_get_layout_plan(slide_data).get("graphic_type") or "")).strip()
    return layout_graphic or "none"


def _resolve_graphic_data(slide_data):
    if isinstance(slide_data.get("graphic_data"), dict):
        return slide_data.get("graphic_data")
    layout_data = _get_layout_plan(slide_data).get("graphic_data")
    return layout_data if isinstance(layout_data, dict) else {}


def has_active_graphic(slide_data):
    graphic_type = _resolve_graphic_type(slide_data)
    graphic_data = _resolve_graphic_data(slide_data)
    if not graphic_type or graphic_type == "none":
        return False
    return isinstance(graphic_data, dict) and bool(graphic_data)


def _graphic_density(slide_data):
    data = _resolve_graphic_data(slide_data)
    if not isinstance(data, dict):
        return 0, 0
    node_count = 0
    text_len = 0
    for key in ("nodes", "cards", "steps", "branches", "events", "rows", "examples", "useful_expressions"):
        value = data.get(key)
        if isinstance(value, list):
            node_count += len(value)
            for item in value:
                if isinstance(item, dict):
                    text_len += len(" ".join(str(v) for v in item.values()))
                else:
                    text_len += len(str(item))
    return node_count, text_len


def safe_compact_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    items = normalize_bullets(slide_data.get("visible_content", []), max_items=MAX_MAIN_BLOCKS, max_chars_each=56)
    if not items:
        items = ["Focus on the key task.", "Use one clear piece of evidence.", "Share your answer briefly."]
    _add_panel(slide, Inches(1.1), Inches(1.95), Inches(10.9), Inches(3.95), style["surface"], line_color=style["soft"])
    _add_bullets(slide, items, Inches(1.45), Inches(2.35), Inches(10.2), Inches(2.35), style, size=12)
    key_sentence = truncate_text(str(slide_data.get("key_sentence") or "").strip(), 110)
    if key_sentence:
        _add_panel(slide, Inches(1.45), Inches(4.95), Inches(10.2), Inches(0.6), style["soft"], line_color=style["soft"])
        _add_textbox(slide, Inches(1.65), Inches(5.13), Inches(9.8), Inches(0.2), f"Key sentence: {key_sentence}", size=11, color=style["primary"], bold=True)
    _add_textbox(slide, Inches(8.6), Inches(6.12), Inches(3.6), Inches(0.2), "More details in teacher script.", size=9, color=style["muted"], align=PP_ALIGN.RIGHT)


def _render_graphic_only_slide(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    key_sentence = truncate_text(str(slide_data.get("key_sentence") or "").strip(), 180)
    if key_sentence:
        _add_panel(slide, Inches(0.95), Inches(1.38), Inches(11.1), Inches(0.32), style["soft"], line_color=style["soft"])
        _add_textbox(
            slide,
            Inches(1.12),
            Inches(1.45),
            Inches(10.7),
            Inches(0.18),
            text=f"Key Sentence: {key_sentence}",
            size=10,
            color=style["primary"],
            bold=True,
        )

    graphic_type = _resolve_graphic_type(slide_data)
    graphic_data = _resolve_graphic_data(slide_data)
    graphic_renderer = GRAPHIC_RENDERERS.get(graphic_type)
    if not graphic_renderer:
        raise ValueError(f"Unsupported graphic_type: {graphic_type}")
    graphic_renderer(slide, prs, graphic_data, area=GRAPHIC_AREA, style=style)

    # Keep image suggestion as a small hint when a full graphic is active.
    image_hint = truncate_text(str(slide_data.get("image_suggestion") or "").strip(), 120)
    if image_hint:
        _add_textbox(
            slide,
            Inches(8.0),
            Inches(6.12),
            Inches(4.2),
            Inches(0.2),
            text=f"Image suggestion: {image_hint}",
            size=8,
            color=style["muted"],
            align=PP_ALIGN.RIGHT,
        )


def render_numbered_cards(slide, slide_data, task, prs, style, template):
    objectives_layout(slide, slide_data, task, prs, style)


def render_question_cards(slide, slide_data, task, prs, style, template):
    leadin_question_layout(slide, slide_data, task, prs, style)


def render_three_step_flow(slide, slide_data, task, prs, style, template):
    prediction_flow_layout(slide, slide_data, task, prs, style)


def render_task_focus_card(slide, slide_data, task, prs, style, template):
    reading_task_layout(slide, slide_data, task, prs, style)


def render_question_evidence_columns(slide, slide_data, task, prs, style, template):
    reading_task_layout(slide, slide_data, task, prs, style)


def render_vocabulary_cards_template(slide, slide_data, task, prs, style, template):
    vocabulary_card_layout(slide, slide_data, task, prs, style)


def render_pattern_practice_blocks(slide, slide_data, task, prs, style, template):
    sentence_analysis_layout(slide, slide_data, task, prs, style)


def render_discussion_prompt_panel(slide, slide_data, task, prs, style, template):
    discussion_layout(slide, slide_data, task, prs, style)


def render_three_takeaway_cards(slide, slide_data, task, prs, style, template):
    mindmap_layout(slide, slide_data, task, prs, style)


def render_homework_levels(slide, slide_data, task, prs, style, template):
    homework_layout(slide, slide_data, task, prs, style)


def render_board_structure(slide, slide_data, task, prs, style, template):
    blackboard_layout(slide, slide_data, task, prs, style)


TEMPLATE_VARIANT_RENDERERS = {
    "numbered_cards": render_numbered_cards,
    "question_cards": render_question_cards,
    "three_step_flow": render_three_step_flow,
    "task_focus_card": render_task_focus_card,
    "question_evidence_columns": render_question_evidence_columns,
    "vocabulary_cards": render_vocabulary_cards_template,
    "pattern_practice_blocks": render_pattern_practice_blocks,
    "discussion_prompt_panel": render_discussion_prompt_panel,
    "three_takeaway_cards": render_three_takeaway_cards,
    "homework_levels": render_homework_levels,
    "board_structure": render_board_structure,
}


def render_by_layout_template(slide, slide_data, task, template, prs, style):
    visual_variant = str(template.get("visual_variant") or "").strip()
    variant_renderer = TEMPLATE_VARIANT_RENDERERS.get(visual_variant)
    if variant_renderer:
        variant_renderer(slide, slide_data, task, prs, style, template)
        return visual_variant or "variant"

    preferred_layout = str(template.get("preferred_layout") or "").strip()
    preferred_renderer = LAYOUT_RENDERERS.get(preferred_layout)
    if preferred_renderer:
        preferred_renderer(slide, slide_data, task, prs, style)
        return preferred_layout

    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    max_blocks = int(template.get("max_blocks") or 3)
    max_items = int(template.get("max_items_per_block") or 2)
    body_size = max(int(template.get("body_font_size") or 18), int(template.get("min_font_size") or 14))
    items = normalize_bullets(slide_data.get("visible_content", []), max_items=max_blocks * max_items, max_chars_each=72)
    if not items:
        items = ["Focus on one key classroom objective."]
    blocks = [items[i:i + max_items] for i in range(0, min(len(items), max_blocks * max_items), max_items)]
    panel_h = Inches(4.25 / max(max_blocks, 1))
    for idx, block in enumerate(blocks[:max_blocks]):
        top = Inches(1.85) + idx * panel_h
        _add_panel(slide, Inches(1.05), top, Inches(10.9), panel_h - Inches(0.08), style["surface"], line_color=style["soft"])
        _add_bullets(
            slide,
            block,
            Inches(1.35),
            top + Inches(0.18),
            Inches(10.2),
            panel_h - Inches(0.3),
            style,
            size=body_size,
        )
    if len(items) > max_blocks * max_items:
        _add_textbox(slide, Inches(8.5), Inches(6.12), Inches(3.7), Inches(0.2), "More details in teacher script.", size=11, color=style["muted"], align=PP_ALIGN.RIGHT)
    return "plain_fallback"


def _add_chinese_hint(slide, slide_data, style):
    hint = str(slide_data.get("chinese_hint") or "").strip()
    if not hint:
        return
    _add_textbox(
        slide,
        Inches(0.95),
        Inches(6.48),
        Inches(9.75),
        Inches(0.18),
        text=f"Hint: {hint}",
        size=8,
        color=style["muted"],
    )


def _draw_image_placeholder(slide, left, top, width, height, style, label):
    box = _add_panel(slide, left, top, width, height, style["soft"], line_color=style["secondary"], rounded=False)
    box.line.width = Pt(1.5)
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        left,
        top,
        left + width,
        top + height,
    ).line.color.rgb = _rgb(style["secondary"])
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        left + width,
        top,
        left,
        top + height,
    ).line.color.rgb = _rgb(style["secondary"])
    _add_textbox(
        slide,
        left + Inches(0.25),
        top + Inches(0.55),
        width - Inches(0.5),
        Inches(0.9),
        text=f"Image suggestion:\n{truncate_text(label, 90)}",
        size=13,
        color=style["secondary"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def _add_tag_strip(slide, items, left, top, width, style, fill=None, text_color=None):
    items = _limit_items(items, limit=4)
    if not items:
        return
    card_width = min(float(width.inches) / max(len(items), 1) - 0.1, 2.6)
    card_width = max(card_width, 1.35)
    for index, item in enumerate(items):
        card_left = left + Inches(index * (card_width + 0.16))
        _add_panel(
            slide,
            card_left,
            top,
            Inches(card_width),
            Inches(0.62),
            fill or style["soft"],
            line_color=style["soft"],
        )
        _add_textbox(
            slide,
            card_left + Inches(0.12),
            top + Inches(0.16),
            Inches(card_width - 0.24),
            Inches(0.22),
            text=item,
            size=10,
            color=text_color or style["text"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def preserve_cover_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_panel(slide, Inches(0.75), Inches(0.8), Inches(5.2), Inches(5.5), style["surface"])
    _add_panel(slide, Inches(6.25), Inches(0.55), Inches(6.05), Inches(5.95), style["soft"], line_color=style["soft"])
    _add_textbox(
        slide,
        Inches(1.05),
        Inches(1.15),
        Inches(4.4),
        Inches(0.38),
        text=task.get("unit") or "Unit",
        size=15,
        color=style["secondary"],
        bold=True,
    )
    _add_textbox(
        slide,
        Inches(1.05),
        Inches(1.75),
        Inches(4.4),
        Inches(1.45),
        text=slide_data["title"],
        size=30,
        color=style["primary"],
        bold=True,
    )
    subtitle = " • ".join(_limit_items(slide_data.get("visible_content", []), limit=3))
    _add_textbox(
        slide,
        Inches(1.05),
        Inches(3.6),
        Inches(4.5),
        Inches(0.8),
        text=subtitle,
        size=16,
        color=style["muted"],
    )
    _draw_image_placeholder(
        slide,
        Inches(6.55),
        Inches(1.0),
        Inches(5.45),
        Inches(4.2),
        style,
        slide_data.get("image_suggestion") or "Image Placeholder: lesson scene / reading / classroom",
    )
    footer_text = " • ".join(part for part in [task.get("grade"), task.get("textbook"), task.get("style")] if part)
    _add_textbox(
        slide,
        Inches(1.05),
        Inches(5.45),
        Inches(4.8),
        Inches(0.28),
        text=footer_text,
        size=11,
        color=style["secondary"],
        bold=True,
    )
    _add_teacher_hint(slide, slide_data, style)


def cover_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_panel(slide, Inches(0.75), Inches(1.0), Inches(8.9), Inches(4.3), style["surface"])
    _add_panel(slide, Inches(9.75), Inches(-0.15), Inches(3.2), Inches(2.15), style["accent"])

    _add_textbox(
        slide,
        Inches(1.15),
        Inches(1.35),
        Inches(6.8),
        Inches(0.4),
        text=task.get("unit") or "Unit",
        size=16,
        color=style["secondary"],
        bold=True,
    )
    _add_textbox(
        slide,
        Inches(1.15),
        Inches(1.85),
        Inches(7.2),
        Inches(1.65),
        text=slide_data["title"],
        size=fit_font_size(slide_data["title"], 7.2, 1.65, 30, 20),
        color=style["primary"],
        bold=True,
    )
    subtitle = " • ".join(
        part for part in [task.get("lesson_type"), task.get("grade"), task.get("textbook")] if part
    )
    _add_textbox(
        slide,
        Inches(1.15),
        Inches(3.75),
        Inches(7.2),
        Inches(0.4),
        text=subtitle,
        size=17,
        color=style["muted"],
    )

    meta_block = _find_block(_get_content_blocks(slide_data), "meta")
    _add_bullets(
        slide,
        normalize_bullets(_block_items(meta_block, fallback=slide_data.get("visible_content", [])), max_items=4, max_chars_each=70),
        Inches(1.15),
        Inches(4.45),
        Inches(7.0),
        Inches(0.9),
        style,
        size=17,
    )

    _add_panel(slide, Inches(9.3), Inches(2.15), Inches(2.7), Inches(1.35), style["primary"])
    _add_textbox(
        slide,
        Inches(9.55),
        Inches(2.35),
        Inches(2.2),
        Inches(0.8),
        text=f"{slide_data.get('estimated_time', '')}\n{slide_data.get('interaction_type', '')}",
        size=12,
        color=(255, 255, 255),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_teacher_hint(slide, slide_data, style)


def objectives_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    raw_blocks = (_get_content_blocks(slide_data) or [
        {"title": f"Objective {idx + 1}", "items": [item]}
        for idx, item in enumerate(_limit_items(slide_data.get("visible_content", []), 3))
    ])[:5]
    blocks = []
    for block in raw_blocks:
        valid_items = [
            item for item in _block_items(block, []) if item and not _is_orphan_number_marker(item)
        ]
        if valid_items:
            blocks.append({"title": _block_title(block, "Objective"), "items": valid_items})
    if not blocks:
        blocks = [{"title": "Objective 1", "items": ["Understand the lesson goal and complete the key task."]}]

    card_positions = [Inches(0.95), Inches(3.55), Inches(6.15), Inches(8.75), Inches(11.35)]
    card_width = Inches(2.25 if len(blocks) >= 4 else 2.95)
    for index, left in enumerate(card_positions):
        if index >= len(blocks):
            break
        block = blocks[index]
        _add_panel(slide, left, Inches(1.85), card_width, Inches(3.55), style["surface"], line_color=style["soft"])
        _add_textbox(
            slide,
            left + Inches(0.22),
            Inches(2.08),
            card_width - Inches(0.35),
            Inches(0.32),
            text=truncate_text(_block_title(block, f"Objective {index + 1}"), 32),
            size=15,
            color=style["primary"],
            bold=True,
        )
        _add_bullets(
            slide,
            normalize_bullets(_block_items(block, ["Understand one clear learning objective."]), max_items=3, max_chars_each=65),
            left + Inches(0.22),
            Inches(2.55),
            card_width - Inches(0.35),
            Inches(2.35),
            style,
            size=14,
        )

    _add_teacher_hint(slide, slide_data, style)


def leadin_question_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    questions = _limit_items(slide_data.get("visible_content", []), limit=3)
    while len(questions) < 3:
        questions.append("Share one idea with your partner.")
    _draw_image_placeholder(
        slide,
        Inches(0.95),
        Inches(1.78),
        Inches(3.15),
        Inches(3.25),
        style,
        slide_data.get("image_suggestion") or "Image Placeholder: warm-up / scene / topic",
    )
    card_positions = [Inches(4.45), Inches(7.15), Inches(9.85)]
    for index, (left, question) in enumerate(zip(card_positions, questions), start=1):
        _add_panel(slide, left, Inches(1.88), Inches(2.15), Inches(2.9), style["surface"], line_color=style["soft"])
        _add_textbox(slide, left + Inches(0.18), Inches(2.1), Inches(1.2), Inches(0.22), f"Q{index}", size=12, color=style["secondary"], bold=True)
        _add_textbox(slide, left + Inches(0.18), Inches(2.5), Inches(1.8), Inches(1.5), question, size=16, color=style["primary"], bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    if slide_data.get("useful_expressions"):
        _add_panel(slide, Inches(0.95), Inches(5.35), Inches(11.15), Inches(0.8), style["soft"], line_color=style["soft"])
        _add_textbox(slide, Inches(1.15), Inches(5.52), Inches(2.0), Inches(0.22), "Useful Expressions", size=11, color=style["primary"], bold=True)
        _add_tag_strip(slide, slide_data.get("useful_expressions", []), Inches(3.2), Inches(5.42), Inches(8.4), style)
    _add_teacher_hint(slide, slide_data, style)


def prediction_flow_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    before_block = _find_block(blocks, "before_reading") or {"title": "Before Reading", "items": slide_data.get("visible_content", [])[:1]}
    predict_block = _find_block(blocks, "predict") or {"title": "Predict", "items": slide_data.get("visible_content", [])[1:3]}
    share_block = _find_block(blocks, "share") or {"title": "Share", "items": slide_data.get("visible_content", [])[3:]}
    expression_block = _find_block(blocks, "expression_support") or {"title": "Useful Expressions", "items": slide_data.get("useful_expressions", [])}
    image_block = _find_block(blocks, "image_placeholder")

    for left, block in [
        (Inches(0.95), before_block),
        (Inches(4.25), predict_block),
        (Inches(7.55), share_block),
    ]:
        _add_panel(slide, left, Inches(1.95), Inches(2.65), Inches(2.25), style["surface"], line_color=style["soft"])
        _add_textbox(slide, left + Inches(0.18), Inches(2.18), Inches(2.0), Inches(0.22), _block_title(block, "Step"), size=14, color=style["primary"], bold=True, align=PP_ALIGN.CENTER)
        _add_bullets(slide, _block_items(block, ["Complete the step."]), left + Inches(0.18), Inches(2.62), Inches(2.0), Inches(1.6), style, size=15)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.6), Inches(3.25), Inches(4.15), Inches(3.25)).line.color.rgb = _rgb(style["accent"])
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.9), Inches(3.25), Inches(7.45), Inches(3.25)).line.color.rgb = _rgb(style["accent"])

    _draw_image_placeholder(
        slide,
        Inches(10.35),
        Inches(1.95),
        Inches(1.8),
        Inches(2.55),
        style,
        _block_items(image_block, [slide_data.get("image_suggestion") or "Image Placeholder: before reading / scene"])[0],
    )
    if _block_items(expression_block, []):
        _add_panel(slide, Inches(1.15), Inches(5.15), Inches(8.9), Inches(0.85), style["soft"], line_color=style["soft"])
        _add_textbox(slide, Inches(1.35), Inches(5.34), Inches(2.1), Inches(0.2), "Useful Expressions", size=11, color=style["primary"], bold=True)
        _add_tag_strip(slide, _block_items(expression_block, []), Inches(3.5), Inches(5.25), Inches(6.1), style)
    _add_teacher_hint(slide, slide_data, style)


def warmup_card_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    _draw_image_placeholder(
        slide,
        Inches(8.95),
        Inches(1.85),
        Inches(3.15),
        Inches(2.45),
        style,
        slide_data.get("image_suggestion") or "Image Placeholder: warm-up / topic",
    )
    _add_panel(slide, Inches(0.95), Inches(1.85), Inches(7.45), Inches(1.25), style["surface"], line_color=style["accent"])
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(2.16),
        Inches(6.95),
        Inches(0.62),
        text=slide_data.get("key_sentence") or slide_data["title"],
        size=20,
        color=style["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    if slide_data.get("possible_answers"):
        _add_textbox(slide, Inches(1.1), Inches(3.45), Inches(2.0), Inches(0.2), "Possible Answers", size=12, color=style["secondary"], bold=True)
        _add_tag_strip(slide, slide_data.get("possible_answers", []), Inches(1.05), Inches(3.72), Inches(7.2), style, fill=style["accent"], text_color=(255, 255, 255))
    if slide_data.get("useful_expressions"):
        _add_panel(slide, Inches(0.95), Inches(4.75), Inches(11.15), Inches(0.95), style["soft"], line_color=style["soft"])
        _add_textbox(slide, Inches(1.2), Inches(4.97), Inches(2.0), Inches(0.2), "Useful Expressions", size=11, color=style["primary"], bold=True)
        _add_tag_strip(slide, slide_data.get("useful_expressions", []), Inches(3.35), Inches(4.88), Inches(8.4), style)
    _add_teacher_hint(slide, slide_data, style)


def image_question_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    image_block = _find_block(blocks, "image_placeholder")
    question_block = _find_block(blocks, "question_panel")

    _draw_image_placeholder(
        slide,
        Inches(0.9),
        Inches(1.65),
        Inches(5.1),
        Inches(3.95),
        style,
        _block_items(image_block, ["Image Placeholder: topic / scene / keyword"])[0],
    )

    _add_panel(slide, Inches(6.35), Inches(1.65), Inches(5.9), Inches(3.95), style["surface"], line_color=style["soft"])
    _add_textbox(
        slide,
        Inches(6.7),
        Inches(1.95),
        Inches(5.0),
        Inches(0.38),
        text=_block_title(question_block, slide_data["title"]),
        size=20,
        color=style["primary"],
        bold=True,
    )
    _add_bullets(
        slide,
        _block_items(question_block, slide_data.get("visible_content", [])),
        Inches(6.7),
        Inches(2.45),
        Inches(5.0),
        Inches(2.25),
        style,
        size=18,
    )
    _add_textbox(
        slide,
        Inches(6.7),
        Inches(5.0),
        Inches(4.9),
        Inches(0.28),
        text=f"Interaction: {slide_data.get('interaction_type', '')}",
        size=10,
        color=style["secondary"],
        bold=True,
    )
    _add_teacher_hint(slide, slide_data, style)


def reading_task_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    primary = _find_block(blocks, "task_steps") or {"title": "Task Steps", "items": slide_data.get("visible_content", [])}
    secondary = _find_block(blocks, "checklist") or {"title": "Do This", "items": []}

    source_items = normalize_bullets(primary.get("items", []), max_items=4, max_chars_each=70)
    cards = [
        (Inches(0.9), source_items[:1] or ["Read quickly and find the main idea."], "Fast Reading"),
        (Inches(4.45), source_items[1:2] or ["Match headings with paragraphs."], "Fast Reading"),
        (Inches(8.0), source_items[2:3] or ["Answer with text evidence."], "Careful Reading"),
    ]
    for left, items, title in cards:
        _add_panel(slide, left, Inches(1.95), Inches(2.95), Inches(3.15), style["surface"], line_color=style["soft"])
        _add_textbox(slide, left + Inches(0.22), Inches(2.18), Inches(2.4), Inches(0.28), title, size=14, color=style["primary"], bold=True)
        _add_bullets(slide, items or ["Find one clear answer with text evidence."], left + Inches(0.22), Inches(2.55), Inches(2.45), Inches(1.8), style, size=15)

    key_sentence = truncate_text(slide_data.get("key_sentence") or "", 130)
    if key_sentence:
        _add_panel(slide, Inches(0.95), Inches(5.05), Inches(10.9), Inches(0.32), style["accent"])
        _add_textbox(slide, Inches(1.15), Inches(5.12), Inches(10.5), Inches(0.18), f"Key Sentence: {key_sentence}", size=10, color=(255, 255, 255), bold=True)

    _add_panel(slide, Inches(0.95), Inches(5.45), Inches(10.1), Inches(0.65), style["soft"])
    checklist_items = _block_items(secondary, ["Check your answers and prepare to share."])
    _add_textbox(
        slide,
        Inches(1.18),
        Inches(5.63),
        Inches(9.5),
        Inches(0.18),
        text="Task Reminder: " + " | ".join(checklist_items),
        size=11,
        color=style["text"],
        bold=True,
    )
    _add_teacher_hint(slide, slide_data, style)


def comparison_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    left_block = _find_block(blocks, "left_column") or {"title": "Focus A", "items": slide_data.get("visible_content", [])[:2]}
    right_block = _find_block(blocks, "right_column") or {"title": "Focus B", "items": slide_data.get("visible_content", [])[2:4]}

    for left, block in [(Inches(0.95), left_block), (Inches(6.55), right_block)]:
        _add_panel(slide, left, Inches(1.85), Inches(5.05), Inches(4.45), style["surface"], line_color=style["soft"])
        _add_textbox(slide, left + Inches(0.25), Inches(2.08), Inches(2.7), Inches(0.32), _block_title(block, "Focus"), size=17, color=style["primary"], bold=True)
        _add_bullets(slide, _block_items(block, ["Add key comparison points."]), left + Inches(0.25), Inches(2.55), Inches(4.35), Inches(2.85), style, size=18)

    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(6.02),
        Inches(3.95),
        Inches(6.47),
        Inches(3.95),
    ).line.color.rgb = _rgb(style["accent"])
    _add_teacher_hint(slide, slide_data, style)


def vocabulary_card_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    positions = [Inches(0.95), Inches(3.95), Inches(6.95), Inches(9.95)]

    for index, left in enumerate(positions):
        block = blocks[index] if index < len(blocks) else {
            "title": f"Word {index + 1}",
            "items": [
                f"Word: keyword_{index + 1}",
                "Meaning: simple classroom meaning",
                "Example: Use it in one sentence.",
            ],
        }
        _add_panel(slide, left, Inches(1.95), Inches(2.25), Inches(3.7), style["surface"], line_color=style["accent"])
        _add_panel(slide, left + Inches(0.18), Inches(2.12), Inches(0.72), Inches(0.38), style["accent"])
        _add_textbox(slide, left + Inches(0.34), Inches(2.18), Inches(0.4), Inches(0.18), str(index + 1), size=12, color=(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
        _add_textbox(slide, left + Inches(1.0), Inches(2.08), Inches(1.1), Inches(0.28), truncate_text(_block_title(block, f"Card {index + 1}"), 18), size=13, color=style["primary"], bold=True)
        _add_bullets(slide, normalize_bullets(_block_items(block, ["Word: term", "Meaning: ...", "Example: ..."]), max_items=3, max_chars_each=52), left + Inches(0.24), Inches(2.72), Inches(1.85), Inches(1.95), style, size=12)

    _add_teacher_hint(slide, slide_data, style)


def sentence_analysis_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    sentence_block = _find_block(blocks, "sentence_box") or {"title": "Sentence Focus", "items": slide_data.get("visible_content", [])[:1]}
    analysis_block = _find_block(blocks, "analysis_steps") or {"title": "Analysis", "items": slide_data.get("visible_content", [])[1:]}

    _add_panel(slide, Inches(0.95), Inches(1.8), Inches(11.2), Inches(1.0), style["surface"], line_color=style["soft"])
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(2.02),
        Inches(10.7),
        Inches(0.45),
        text=_block_items(sentence_block, ["Add the key sentence."])[0],
        size=20,
        color=style["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    columns = _chunk_analysis_items(_block_items(analysis_block, ["Mark the key part.", "Explain the meaning.", "Use it in context."]))
    for index, (left, item) in enumerate(zip([Inches(1.05), Inches(4.35), Inches(7.65)], columns), start=1):
        _add_panel(slide, left, Inches(3.2), Inches(2.75), Inches(2.3), style["soft"])
        _add_textbox(slide, left + Inches(0.22), Inches(3.4), Inches(2.25), Inches(0.22), f"Point {index}", size=13, color=style["secondary"], bold=True)
        _add_bullets(slide, [item], left + Inches(0.22), Inches(3.75), Inches(2.15), Inches(1.2), style, size=16)

    if slide_data.get("useful_expressions"):
        _add_panel(slide, Inches(0.95), Inches(5.15), Inches(6.1), Inches(0.82), style["soft"], line_color=style["soft"])
        _add_textbox(slide, Inches(1.15), Inches(5.34), Inches(1.9), Inches(0.2), "Useful Expressions", size=11, color=style["primary"], bold=True)
        _add_tag_strip(slide, slide_data.get("useful_expressions", []), Inches(3.0), Inches(5.24), Inches(3.8), style)
    if slide_data.get("possible_answers"):
        _add_panel(slide, Inches(7.25), Inches(5.15), Inches(4.8), Inches(0.82), style["accent"], line_color=style["accent"])
        _add_textbox(slide, Inches(7.5), Inches(5.34), Inches(1.6), Inches(0.2), "Possible Answers", size=11, color=(255, 255, 255), bold=True)
        _add_tag_strip(slide, slide_data.get("possible_answers", []), Inches(8.95), Inches(5.24), Inches(2.7), style, fill=(255, 244, 221), text_color=style["primary"])

    _add_teacher_hint(slide, slide_data, style)


def _chunk_analysis_items(items):
    items = items[:3]
    if len(items) < 3:
        items += ["Add one more focus point."] * (3 - len(items))
    return items


def discussion_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    prompt_block = _find_block(blocks, "discussion_prompt") or {"title": "Topic", "items": slide_data.get("visible_content", [])[:1] or ["Discuss the key idea of the lesson."]}
    roles_block = _find_block(blocks, "group_roles") or {"title": "Group Task", "items": slide_data.get("visible_content", [])[1:3] or ["List two supporting points.", "Decide one speaker."]}
    output_block = _find_block(blocks, "share_out") or {"title": "Share", "items": slide_data.get("visible_content", [])[3:] or ["Share one clear group conclusion."]}
    support_block = _find_block(blocks, "expression_support") or {"title": "Useful Expressions", "items": slide_data.get("useful_expressions", [])}

    center = _add_panel(slide, Inches(4.7), Inches(2.3), Inches(3.6), Inches(2.0), style["surface"], line_color=style["accent"])
    center.line.width = Pt(2)
    _add_textbox(slide, Inches(5.0), Inches(2.55), Inches(3.0), Inches(0.32), _block_title(prompt_block, "Discuss"), size=17, color=style["primary"], bold=True, align=PP_ALIGN.CENTER)
    _add_bullets(slide, _block_items(prompt_block, ["Share one central idea."]), Inches(5.1), Inches(3.0), Inches(2.8), Inches(0.9), style, size=16)

    for left, top, block in [
        (Inches(1.1), Inches(2.55), roles_block),
        (Inches(8.9), Inches(2.55), output_block),
    ]:
        _add_panel(slide, left, top, Inches(2.9), Inches(1.8), style["soft"])
        _add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.3), Inches(0.24), _block_title(block, "Task"), size=14, color=style["primary"], bold=True)
        _add_bullets(slide, _block_items(block, ["Prepare one class response."]), left + Inches(0.2), top + Inches(0.55), Inches(2.4), Inches(0.9), style, size=15)

    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.0), Inches(3.45), Inches(4.65), Inches(3.45)).line.color.rgb = _rgb(style["accent"])
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.32), Inches(3.45), Inches(8.85), Inches(3.45)).line.color.rgb = _rgb(style["accent"])
    support_items = normalize_bullets(_block_items(support_block, []), max_items=4, max_chars_each=42)
    if support_items:
        _add_panel(slide, Inches(2.25), Inches(5.1), Inches(8.6), Inches(0.82), style["soft"], line_color=style["soft"])
        _add_textbox(slide, Inches(2.5), Inches(5.3), Inches(1.8), Inches(0.2), "Useful Expressions", size=11, color=style["primary"], bold=True)
        _add_tag_strip(slide, support_items, Inches(4.35), Inches(5.2), Inches(5.9), style)
    _add_teacher_hint(slide, slide_data, style)


def mindmap_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    center_block = _find_block(blocks, "center_topic") or {"title": "Center", "items": slide_data.get("visible_content", [])[:1]}
    branches_block = _find_block(blocks, "branches") or {"title": "Branches", "items": slide_data.get("visible_content", [])[1:]}
    branches = normalize_bullets(_block_items(branches_block, ["What we learned", "How we learned", "What to do next"]), max_items=4, max_chars_each=40)

    center = _add_panel(slide, Inches(4.8), Inches(2.55), Inches(3.2), Inches(1.4), style["surface"], line_color=style["accent"])
    center.line.width = Pt(2)
    _add_textbox(slide, Inches(5.1), Inches(2.95), Inches(2.6), Inches(0.45), _block_items(center_block, [slide_data["title"]])[0], size=18, color=style["primary"], bold=True, align=PP_ALIGN.CENTER)

    positions = [
        (Inches(1.15), Inches(1.8)),
        (Inches(8.95), Inches(1.8)),
        (Inches(1.15), Inches(4.45)),
        (Inches(8.95), Inches(4.45)),
    ]
    for (left, top), item in zip(positions, branches):
        _add_panel(slide, left, top, Inches(2.8), Inches(0.92), style["soft"])
        _add_textbox(slide, left + Inches(0.18), top + Inches(0.25), Inches(2.4), Inches(0.25), item, size=14, color=style["text"], bold=True, align=PP_ALIGN.CENTER)
        start_x = left + Inches(2.8 if left < Inches(5) else 0)
        end_x = Inches(4.8 if left < Inches(5) else 8.0)
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, top + Inches(0.46), end_x, Inches(3.25)).line.color.rgb = _rgb(style["accent"])

    _add_teacher_hint(slide, slide_data, style)


def homework_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    required_block = _find_block(blocks, "required_task") or {"title": "Basic", "items": slide_data.get("visible_content", [])[:2] or ["Review target words and one key sentence.", "Finish the guided task from class."]}
    extension_block = _find_block(blocks, "extension_task") or {"title": "Optional / Challenge", "items": slide_data.get("visible_content", [])[2:4] or ["Write one extra paragraph with evidence.", "Prepare one question for next lesson."]}
    reminder_block = _find_block(blocks, "reminder") or {"title": "Reminder", "items": slide_data.get("visible_content", [])[4:]}

    for left, block, fill in [
        (Inches(0.95), required_block, style["surface"]),
        (Inches(6.65), extension_block, style["soft"]),
    ]:
        _add_panel(slide, left, Inches(1.85), Inches(5.1), Inches(3.2), fill, line_color=style["soft"])
        _add_textbox(slide, left + Inches(0.22), Inches(2.12), Inches(2.8), Inches(0.3), _block_title(block, "Task"), size=17, color=style["primary"], bold=True)
        _add_bullets(slide, normalize_bullets(_block_items(block, ["Complete one specific homework product with clear quality target."]), max_items=2, max_chars_each=70), left + Inches(0.22), Inches(2.58), Inches(4.5), Inches(2.15), style, size=15)

    _add_panel(slide, Inches(1.45), Inches(5.45), Inches(10.3), Inches(0.65), style["accent"])
    _add_textbox(
        slide,
        Inches(1.7),
        Inches(5.64),
        Inches(9.8),
        Inches(0.22),
        text="Reminder: " + " | ".join(_block_items(reminder_block, ["Bring one question to the next class."])),
        size=11,
        color=(255, 255, 255),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_teacher_hint(slide, slide_data, style)


def blackboard_layout(slide, slide_data, task, prs, style):
    _set_background(slide, prs, style)
    _add_header_band(slide, slide_data, style)
    blocks = _get_content_blocks(slide_data)
    header_block = _find_block(blocks, "board_header") or {"title": "Board Theme", "items": slide_data.get("visible_content", [])[:1]}
    left_block = _find_block(blocks, "board_left") or {"title": "Language", "items": slide_data.get("visible_content", [])[1:3]}
    right_block = _find_block(blocks, "board_right") or {"title": "Task", "items": slide_data.get("visible_content", [])[3:]}

    board = _add_panel(slide, Inches(0.95), Inches(1.72), Inches(11.3), Inches(4.55), (40, 73, 52), line_color=(28, 50, 36), rounded=False)
    board.line.width = Pt(2)
    _add_textbox(slide, Inches(1.25), Inches(1.98), Inches(10.6), Inches(0.35), _block_items(header_block, [slide_data["title"]])[0], size=22, color=(247, 250, 244), bold=True, align=PP_ALIGN.CENTER)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.6), Inches(2.55), Inches(6.6), Inches(5.8)).line.color.rgb = RGBColor(225, 235, 224)
    _add_textbox(slide, Inches(1.28), Inches(2.7), Inches(4.8), Inches(0.3), _block_title(left_block, "Language"), size=16, color=(233, 241, 228), bold=True)
    _add_textbox(slide, Inches(6.9), Inches(2.7), Inches(4.6), Inches(0.3), _block_title(right_block, "Task"), size=16, color=(233, 241, 228), bold=True)
    _add_bullets(slide, _block_items(left_block, ["Key language", "Useful frame"]), Inches(1.28), Inches(3.12), Inches(4.8), Inches(2.0), {"text": (245, 247, 242)}, size=17)
    _add_bullets(slide, _block_items(right_block, ["Task target", "Class output"]), Inches(6.9), Inches(3.12), Inches(4.2), Inches(2.0), {"text": (245, 247, 242)}, size=17)
    _add_teacher_hint(slide, slide_data, style)


LAYOUT_RENDERERS = {
    "cover_layout": cover_layout,
    "preserve_cover_layout": preserve_cover_layout,
    "objectives_layout": objectives_layout,
    "leadin_question_layout": leadin_question_layout,
    "prediction_flow_layout": prediction_flow_layout,
    "warmup_card_layout": warmup_card_layout,
    "image_question_layout": image_question_layout,
    "reading_task_layout": reading_task_layout,
    "comparison_layout": comparison_layout,
    "vocabulary_card_layout": vocabulary_card_layout,
    "sentence_analysis_layout": sentence_analysis_layout,
    "discussion_layout": discussion_layout,
    "mindmap_layout": mindmap_layout,
    "homework_layout": homework_layout,
    "blackboard_layout": blackboard_layout,
}


def validate_pptx_file(path):
    """Run a lightweight structural check on the exported pptx file."""

    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {path}")
    if path.stat().st_size <= 0:
        raise ValueError(f"PPTX file is empty: {path}")
    if not zipfile.is_zipfile(path):
        raise ValueError(f"PPTX file is not a valid zip package: {path}")

    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        required_entries = {
            "ppt/presentation.xml",
            "[Content_Types].xml",
            "ppt/slides/slide1.xml",
        }
        missing = required_entries - names
        if missing:
            raise ValueError(f"PPTX package is missing required entries: {sorted(missing)}")

    return True


def export_pptx(task, slides):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_SIZE[0])
    prs.slide_height = Inches(SLIDE_SIZE[1])
    default_style = get_ppt_style_profile(task)

    # To maximize compatibility with PowerPoint, Keynote, and WPS,
    # this export intentionally does not write speaker notes / notes slides.
    # teacher_notes remain available in the editor, JSON preview, and DOCX script export.
    for slide_data in slides:
        render_data = dict(slide_data)
        if not render_data.get("layout_plan"):
            render_data["layout_plan"] = plan_layout_for_slide(render_data, task)
        slide_task = dict(task)
        if render_data.get("ppt_style"):
            slide_task["ppt_style"] = render_data.get("ppt_style")
        style = get_ppt_style_profile(slide_task) if slide_task.get("ppt_style") else default_style
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        layout_template = (_get_layout_plan(render_data).get("layout_template") or
                           get_template_for_slide(render_data.get("slide_type")))
        layout_type = (_get_layout_plan(render_data).get("layout_type") or "reading_task_layout")
        renderer = LAYOUT_RENDERERS.get(layout_type, reading_task_layout)
        graphic_active = has_active_graphic(render_data)
        node_count, text_len = _graphic_density(render_data)
        dense_graphic = node_count > 4 or text_len > 280
        dense_slide = len(render_data.get("visible_content", []) or []) > 4
        compact_risky_type = str(render_data.get("slide_type") or "") in {"lead_in", "group_discussion", "language_points"}
        force_compact = dense_graphic or (compact_risky_type and dense_slide)
        if layout_template:
            rendered_variant = render_by_layout_template(slide, render_data, task, layout_template, prs, style)
            render_data.setdefault("layout_plan", {})["rendered_visual_variant"] = rendered_variant
        elif graphic_active:
            try:
                if force_compact:
                    safe_compact_layout(slide, render_data, task, prs, style)
                else:
                    _render_graphic_only_slide(slide, render_data, task, prs, style)
            except Exception:
                # Full fallback: if graphic pipeline fails, render legacy layout only.
                safe_compact_layout(slide, render_data, task, prs, style)
        else:
            if force_compact:
                safe_compact_layout(slide, render_data, task, prs, style)
            else:
                renderer(slide, render_data, task, prs, style)
        _add_chinese_hint(slide, render_data, style)
        _add_footer(slide, render_data.get("slide_index", 1), task, style)

    filename = safe_course_filename(task, "slides.pptx")
    path = EXPORT_PPTX_DIR / filename
    prs.save(path)
    validate_pptx_file(path)
    return path
