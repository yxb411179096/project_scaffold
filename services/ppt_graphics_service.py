"""Rule-based teaching graphics renderer for PPTX.

All rendering uses python-pptx public APIs only.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

GRAPHIC_TITLE_FONT = 18
GRAPHIC_CARD_TITLE_FONT = 14
GRAPHIC_BODY_FONT = 12
GRAPHIC_SMALL_FONT = 11
GRAPHIC_MIN_FONT = 11


def _rgb(color):
    return RGBColor(*color)


def _panel(slide, left, top, width, height, fill, border=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if border:
        shape.line.color.rgb = _rgb(border)
    else:
        shape.line.fill.background()
    return shape


def _safe_font(size):
    return max(GRAPHIC_MIN_FONT, int(size))


def _text(slide, left, top, width, height, value, size=12, color=(32, 43, 57), bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.text = smart_clip(value, 120)
    p = tf.paragraphs[0]
    p.alignment = align
    p.font.size = Pt(_safe_font(size))
    p.font.color.rgb = _rgb(color)
    p.font.bold = bold
    return box


def smart_clip(value, max_chars=80, preserve_words=True):
    text = str(value or "").strip()
    if len(text) <= max_chars:
        return text
    if max_chars <= 3:
        return "..."
    target = text[: max_chars - 3].rstrip()
    if preserve_words and " " in target:
        ascii_ratio = sum(1 for c in target if ord(c) < 128) / max(len(target), 1)
        if ascii_ratio > 0.7:
            last_space = target.rfind(" ")
            if last_space > max(5, int(max_chars * 0.5)):
                target = target[:last_space].rstrip()
    return target + "..."


def _split_area(area):
    left, top, width, height = area
    return Inches(left), Inches(top), Inches(width), Inches(height)


def render_mindmap(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    center = str(data.get("center") or data.get("title") or "Summary")
    branches = [smart_clip(x, 32, preserve_words=True) for x in list(data.get("branches") or [])[:4]] or ["What we learned", "Key words", "Skill", "Reflection"]
    center_w = width * 0.34
    center_h = height * 0.34
    center_l = left + (width - center_w) / 2
    center_t = top + (height - center_h) / 2
    _panel(slide, center_l, center_t, center_w, center_h, style["surface"], style["accent"])
    _text(slide, center_l + Inches(0.1), center_t + Inches(0.2), center_w - Inches(0.2), center_h - Inches(0.3), smart_clip(center, 52), GRAPHIC_CARD_TITLE_FONT, style["primary"], True, PP_ALIGN.CENTER)
    slots = [(0.05, 0.1), (0.72, 0.1), (0.05, 0.72), (0.72, 0.72), (0.4, 0.02)]
    for idx, item in enumerate(branches[: len(slots)]):
        x, y = slots[idx]
        bw = width * 0.22
        bh = height * 0.2
        bl = left + width * x
        bt = top + height * y
        _panel(slide, bl, bt, bw, bh, style["soft"], style["soft"])
        _text(slide, bl + Inches(0.08), bt + Inches(0.12), bw - Inches(0.16), bh - Inches(0.2), item, GRAPHIC_BODY_FONT, style["text"], True, PP_ALIGN.CENTER)
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_l + center_w / 2, center_t + center_h / 2, bl + bw / 2, bt + bh / 2).line.color.rgb = _rgb(style["accent"])
    if len(list(data.get("branches") or [])) > 4:
        _text(slide, left + Inches(0.05), top + height - Inches(0.22), Inches(2.6), Inches(0.18), "More details in teacher script.", GRAPHIC_SMALL_FONT, style["muted"])


def render_flowchart(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    steps = [smart_clip(x, 60) for x in list(data.get("steps") or [])[:3]] or ["Before reading", "While reading", "After reading"]
    count = len(steps)
    card_w = (width - Inches(0.2 * (count - 1))) / count
    for idx, step in enumerate(steps):
        cl = left + idx * (card_w + Inches(0.2))
        _panel(slide, cl, top + Inches(0.35), card_w, height - Inches(0.7), style["surface"], style["soft"])
        _text(slide, cl + Inches(0.1), top + Inches(0.7), card_w - Inches(0.2), height - Inches(1.1), step, GRAPHIC_BODY_FONT, style["primary"], True, PP_ALIGN.CENTER)
        if idx < count - 1:
            slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cl + card_w, top + height / 2, cl + card_w + Inches(0.18), top + height / 2).line.color.rgb = _rgb(style["accent"])


def render_timeline(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    events = [smart_clip(x, 28) for x in list(data.get("events") or [])[:4]] or ["Start", "Develop", "Change", "Result"]
    y = top + height / 2
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(0.2), y, left + width - Inches(0.2), y).line.color.rgb = _rgb(style["secondary"])
    step = (width - Inches(0.4)) / max(len(events), 1)
    for idx, event in enumerate(events):
        x = left + Inches(0.2) + idx * step
        _panel(slide, x, y - Inches(0.95), Inches(1.6), Inches(0.7), style["soft"], style["soft"])
        _text(slide, x + Inches(0.08), y - Inches(0.84), Inches(1.44), Inches(0.5), event, GRAPHIC_BODY_FONT, style["text"], True, PP_ALIGN.CENTER)


def render_comparison_table(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    rows = list(data.get("rows") or [])[:3]
    if not rows:
        rows = [{"left": "Point A", "right": "Point B"}, {"left": "Example A", "right": "Example B"}]
    _panel(slide, left, top, width, height, style["surface"], style["soft"])
    mid = left + width / 2
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, mid, top + Inches(0.4), mid, top + height - Inches(0.2)).line.color.rgb = _rgb(style["soft"])
    _text(slide, left + Inches(0.1), top + Inches(0.1), width / 2 - Inches(0.2), Inches(0.25), "A", GRAPHIC_BODY_FONT, style["primary"], True, PP_ALIGN.CENTER)
    _text(slide, mid + Inches(0.1), top + Inches(0.1), width / 2 - Inches(0.2), Inches(0.25), "B", GRAPHIC_BODY_FONT, style["primary"], True, PP_ALIGN.CENTER)
    row_h = (height - Inches(0.55)) / max(len(rows), 1)
    for idx, row in enumerate(rows):
        rt = top + Inches(0.4) + idx * row_h
        _text(slide, left + Inches(0.12), rt + Inches(0.08), width / 2 - Inches(0.24), row_h - Inches(0.16), smart_clip(row.get("left", ""), 56), GRAPHIC_BODY_FONT, style["text"])
        _text(slide, mid + Inches(0.12), rt + Inches(0.08), width / 2 - Inches(0.24), row_h - Inches(0.16), smart_clip(row.get("right", ""), 56), GRAPHIC_BODY_FONT, style["text"])


def render_reading_structure(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    nodes = list(data.get("nodes") or [])[:3]
    if not nodes:
        nodes = [{"label": "Problem", "content": "Identify the problem."}, {"label": "Solution", "content": "Find the solution."}, {"label": "Result", "content": "Summarize the result."}]
    card_w = (width - Inches(0.4)) / max(len(nodes), 1)
    for idx, node in enumerate(nodes):
        cl = left + idx * (card_w + Inches(0.2))
        _panel(slide, cl, top + Inches(0.35), card_w, height - Inches(0.7), style["surface"], style["accent"])
        _text(slide, cl + Inches(0.08), top + Inches(0.5), card_w - Inches(0.16), Inches(0.25), smart_clip(node.get("label", f"Node {idx + 1}"), 22), GRAPHIC_CARD_TITLE_FONT, style["secondary"], True, PP_ALIGN.CENTER)
        _text(slide, cl + Inches(0.08), top + Inches(0.82), card_w - Inches(0.16), height - Inches(1.0), smart_clip(node.get("content", ""), 60), GRAPHIC_BODY_FONT, style["text"])


def render_writing_framework(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    sections = dict(data.get("sections") or {})
    opening = sections.get("Opening", "Opening")
    body = sections.get("Body", "Body")
    ending = sections.get("Ending", "Ending")
    _panel(slide, left, top, width * 0.62, height, style["surface"], style["soft"])
    _text(slide, left + Inches(0.1), top + Inches(0.12), width * 0.58, Inches(0.22), "Writing Structure", GRAPHIC_CARD_TITLE_FONT, style["primary"], True)
    for idx, (label, content) in enumerate([("Opening", opening), ("Body", body), ("Ending", ending)]):
        y = top + Inches(0.45) + idx * Inches(0.75)
        _panel(slide, left + Inches(0.12), y, width * 0.56, Inches(0.62), style["soft"], style["soft"])
        _text(slide, left + Inches(0.2), y + Inches(0.08), Inches(1.0), Inches(0.2), label, GRAPHIC_BODY_FONT, style["secondary"], True)
        _text(slide, left + Inches(1.2), y + Inches(0.08), width * 0.44, Inches(0.45), smart_clip(content, 72), GRAPHIC_BODY_FONT, style["text"])
    right_l = left + width * 0.66
    _panel(slide, right_l, top, width * 0.34, height, style["soft"], style["soft"])
    _text(slide, right_l + Inches(0.08), top + Inches(0.1), width * 0.30, Inches(0.2), "Checklist", GRAPHIC_BODY_FONT, style["primary"], True)
    checklist = list(data.get("checklist") or ["Clear purpose", "Logical structure", "Accurate language"])[:3]
    for idx, item in enumerate(checklist):
        _text(slide, right_l + Inches(0.1), top + Inches(0.35) + idx * Inches(0.45), width * 0.28, Inches(0.35), f"- {smart_clip(item, 34)}", GRAPHIC_SMALL_FONT, style["text"])


def render_grammar_rule_chart(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    _panel(slide, left, top, width, height, style["surface"], style["soft"])
    examples = [smart_clip(x, 55) for x in list(data.get("examples") or ["Example 1", "Example 2"])[:2]]
    rule = smart_clip(data.get("rule") or "Rule summary", 72)
    practice = smart_clip(data.get("practice") or "Practice task", 72)
    mistakes = [smart_clip(x, 40) for x in list(data.get("common_mistakes") or ["Common mistake"])[:1]]
    _text(slide, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.2), "Examples", GRAPHIC_BODY_FONT, style["secondary"], True)
    _text(slide, left + Inches(0.12), top + Inches(0.35), width - Inches(0.24), Inches(0.55), " • " + "  • ".join(examples), GRAPHIC_SMALL_FONT, style["text"])
    _text(slide, left + Inches(0.1), top + Inches(0.98), width - Inches(0.2), Inches(0.2), "Rule", GRAPHIC_BODY_FONT, style["secondary"], True)
    _text(slide, left + Inches(0.12), top + Inches(1.18), width - Inches(0.24), Inches(0.45), rule, GRAPHIC_BODY_FONT, style["text"])
    _text(slide, left + Inches(0.1), top + Inches(1.66), width - Inches(0.2), Inches(0.2), "Practice", GRAPHIC_BODY_FONT, style["secondary"], True)
    _text(slide, left + Inches(0.12), top + Inches(1.86), width - Inches(0.24), Inches(0.45), practice, GRAPHIC_BODY_FONT, style["text"])
    _text(slide, left + width - Inches(3.2), top + height - Inches(0.22), Inches(3.1), Inches(0.18), "More details in teacher script.", GRAPHIC_SMALL_FONT, style["muted"], align=PP_ALIGN.RIGHT)


def render_vocabulary_cards(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    cards = list(data.get("cards") or [])[:3]
    if not cards:
        cards = [{"word": "word", "meaning": "meaning", "example": "example sentence"} for _ in range(3)]
    cards = cards[:3]
    cols = len(cards)
    rows = 1
    card_w = (width - Inches(0.2)) / max(cols, 1)
    card_h = (height - Inches(0.2)) / max(rows, 1)
    for idx, card in enumerate(cards):
        r = idx // cols
        c = idx % cols
        cl = left + c * (card_w + Inches(0.2))
        ct = top + r * (card_h + Inches(0.2))
        word_raw = str(card.get("word", "")).strip()
        meaning_raw = str(card.get("meaning", "")).strip()
        example_raw = str(card.get("example", "")).strip()
        if len(word_raw) > 30:
            meaning_raw = word_raw if not meaning_raw else meaning_raw
            word_raw = f"Item {idx + 1}"
        _panel(slide, cl, ct, card_w, card_h, style["surface"], style["soft"])
        _text(slide, cl + Inches(0.08), ct + Inches(0.08), card_w - Inches(0.16), Inches(0.2), smart_clip(word_raw or f"Item {idx + 1}", 24), GRAPHIC_CARD_TITLE_FONT, style["primary"], True)
        _text(slide, cl + Inches(0.08), ct + Inches(0.34), card_w - Inches(0.16), Inches(0.24), f"Meaning: {smart_clip(meaning_raw, 44)}", GRAPHIC_BODY_FONT, style["text"])
        _text(slide, cl + Inches(0.08), ct + Inches(0.62), card_w - Inches(0.16), card_h - Inches(0.7), f"Example: {smart_clip(example_raw, 50)}", GRAPHIC_SMALL_FONT, style["text"])
    if len(list(data.get("cards") or [])) > 3:
        _text(slide, left + width - Inches(2.8), top + height - Inches(0.22), Inches(2.7), Inches(0.18), "More details in teacher script.", GRAPHIC_SMALL_FONT, style["muted"], align=PP_ALIGN.RIGHT)


def render_task_steps(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    steps = [smart_clip(x, 18) for x in list(data.get("steps") or [])[:3]] or ["Read key task", "Find clues", "Share evidence"]
    count = len(steps)
    card_w = (width - Inches(0.2 * (count - 1))) / count
    for idx, step in enumerate(steps):
        cl = left + idx * (card_w + Inches(0.2))
        _panel(slide, cl, top + Inches(0.25), card_w, height - Inches(0.5), style["soft"], style["soft"])
        _text(slide, cl + Inches(0.08), top + Inches(0.34), card_w - Inches(0.16), Inches(0.22), f"Step {idx + 1}", GRAPHIC_BODY_FONT, style["secondary"], True, PP_ALIGN.CENTER)
        _text(slide, cl + Inches(0.08), top + Inches(0.62), card_w - Inches(0.16), height - Inches(0.9), step, GRAPHIC_BODY_FONT, style["text"])


def render_discussion_grid(slide, prs, data, area, style):
    left, top, width, height = _split_area(area)
    _panel(slide, left, top, width, height, style["surface"], style["soft"])
    topic = smart_clip(data.get("topic") or "Discussion topic", 60)
    useful = [smart_clip(x, 40) for x in list(data.get("useful_expressions") or ["I agree that...", "In my opinion..."])[:3]]
    share = smart_clip(data.get("share_task") or "Share one group conclusion.", 56)
    mid_x = left + width / 2
    mid_y = top + height / 2
    _panel(slide, left + Inches(0.2), top + Inches(0.3), width - Inches(0.4), Inches(0.95), style["soft"], style["soft"])
    _text(slide, left + Inches(0.35), top + Inches(0.42), width - Inches(0.7), Inches(0.6), f"Discussion Topic: {topic}", GRAPHIC_BODY_FONT, style["primary"], True)
    _panel(slide, left + Inches(0.2), top + Inches(1.45), width - Inches(0.4), Inches(1.05), style["surface"], style["soft"])
    _text(slide, left + Inches(0.35), top + Inches(1.58), width - Inches(0.7), Inches(0.8), "Useful expressions:\n- " + "\n- ".join(useful[:2]), GRAPHIC_BODY_FONT, style["text"])
    _panel(slide, left + Inches(0.2), top + Inches(2.65), width - Inches(0.4), Inches(0.95), style["surface"], style["soft"])
    _text(slide, left + Inches(0.35), top + Inches(2.78), width - Inches(0.7), Inches(0.6), f"Pair work task: {share}", GRAPHIC_BODY_FONT, style["text"], True)


GRAPHIC_RENDERERS = {
    "mindmap": render_mindmap,
    "flowchart": render_flowchart,
    "timeline": render_timeline,
    "comparison_table": render_comparison_table,
    "reading_structure": render_reading_structure,
    "writing_framework": render_writing_framework,
    "grammar_rule_chart": render_grammar_rule_chart,
    "vocabulary_cards": render_vocabulary_cards,
    "task_steps": render_task_steps,
    "discussion_grid": render_discussion_grid,
}
