from pptx import Presentation

from services.agents.layout_planner_agent import plan_layout_for_slide
from services.ppt_graphics_service import GRAPHIC_RENDERERS
from services.ppt_quality_check_service import check_ppt_quality
import services.ppt_render_service as ppt_render_service
from services.ppt_render_service import export_pptx, has_active_graphic


def _task(lesson_type="Reading"):
    return {
        "id": 19001,
        "course_title": "[TEST] Round19 Graphics",
        "topic": "The Internet",
        "grade": "高一",
        "textbook": "人教版",
        "volume": "必修二",
        "unit": "Unit 3",
        "lesson_type": lesson_type,
        "duration": 45,
        "student_level": "中等",
        "style": "常规课",
    }


def run():
    # 1/2/3 layout planner graphic assignments
    reading_slide = {"slide_type": "careful_reading", "title": "Careful Reading", "visible_content": ["Find details", "Problem", "Solution"]}
    writing_slide = {"slide_type": "structure_analysis", "title": "Structure", "visible_content": ["Opening", "Body", "Ending"]}
    grammar_slide = {"slide_type": "grammar_rule", "title": "Rule", "visible_content": ["Example", "Rule", "Practice"]}
    leadin_plan = plan_layout_for_slide({"slide_type": "lead_in", "title": "Lead-in", "visible_content": ["Q1", "Q2"]}, _task("Reading"))
    assert leadin_plan.get("graphic_type") != "task_steps"
    assert plan_layout_for_slide(reading_slide, _task("Reading")).get("graphic_type") == "reading_structure"
    assert plan_layout_for_slide(writing_slide, _task("Writing")).get("graphic_type") == "writing_framework"
    assert plan_layout_for_slide(grammar_slide, _task("Grammar")).get("graphic_type") == "grammar_rule_chart"

    # 4 graphics renderers callable without crash
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = {
        "primary": (21, 76, 121),
        "secondary": (51, 115, 167),
        "accent": (244, 178, 102),
        "text": (32, 43, 57),
        "muted": (102, 112, 133),
        "surface": (255, 255, 255),
        "soft": (240, 246, 252),
    }
    for name, renderer in GRAPHIC_RENDERERS.items():
        renderer(slide, prs, {"title": name}, area=(0.8, 1.2, 10.8, 4.8), style=style)

    # 5 task_steps generic detection
    bad_report = check_ppt_quality(
        _task("Reading"),
        [
            {
                "slide_index": 1,
                "slide_type": "lead_in",
                "title": "Lead in",
                "visible_content": ["Step 1", "Step 2", "Step 3"],
                "teacher_notes": "notes",
                "graphic_type": "task_steps",
                "graphic_data": {"title": "Task Steps", "steps": ["Step 1", "Step 2", "Step 3"]},
                "layout_plan": {"layout_type": "reading_task_layout", "graphic_type": "task_steps"},
            }
        ],
    )
    issue_codes = {
        detail["code"]
        for slide_report in bad_report["slides"]
        for detail in slide_report.get("issue_details", [])
    }
    assert "task_steps_too_generic" in issue_codes

    # 6 invalid graphic_type fallback should not fail export
    bad_graphic_slide = {
        "slide_index": 1,
        "slide_type": "summary",
        "title": "Summary",
        "visible_content": ["What we learned", "How we learned", "What to do next"],
        "teacher_notes": "notes",
        "teaching_purpose": "summary",
        "estimated_time": "4 min",
        "interaction_type": "Teacher-student interaction",
        "graphic_type": "invalid_type",
        "layout_plan": {"layout_type": "mindmap_layout"},
    }
    path_bad = export_pptx(_task("Reading"), [bad_graphic_slide])
    assert path_bad.exists()

    # 6.1 graphic active should skip normal body renderer path (mutual exclusion)
    orig_renderer = ppt_render_service.LAYOUT_RENDERERS["reading_task_layout"]
    called = {"renderer": 0}

    def _counting_renderer(slide, slide_data, task, prs, style):
        called["renderer"] += 1
        return orig_renderer(slide, slide_data, task, prs, style)

    ppt_render_service.LAYOUT_RENDERERS["reading_task_layout"] = _counting_renderer
    try:
        graphic_slide = {
            "slide_index": 1,
            "slide_type": "fast_reading",
            "title": "Fast Reading",
            "visible_content": ["Read quickly", "Find main idea", "Share answer", "Check details"],
            "teacher_notes": "notes",
            "teaching_purpose": "reading task",
            "estimated_time": "5 min",
            "interaction_type": "Individual work",
            "graphic_type": "task_steps",
            "graphic_data": {"title": "Task Steps", "steps": ["Read quickly and find the main idea.", "Find one key clue.", "Share your answer with evidence."]},
            "layout_plan": {"layout_type": "reading_task_layout", "graphic_type": "task_steps", "graphic_data": {"title": "Task Steps", "steps": ["A", "B", "C"]}},
        }
        assert has_active_graphic(graphic_slide) is True
        path_graphic = export_pptx(_task("Reading"), [graphic_slide])
        assert path_graphic.exists()
        assert called["renderer"] == 0
    finally:
        ppt_render_service.LAYOUT_RENDERERS["reading_task_layout"] = orig_renderer

    # 7 quality check graphic_data_empty
    report = check_ppt_quality(
        _task("Reading"),
        [
            {
                "slide_index": 1,
                "slide_type": "careful_reading",
                "title": "Careful Reading",
                "visible_content": ["Find evidence", "Explain details"],
                "teacher_notes": "notes",
                "graphic_type": "reading_structure",
                "graphic_data": {},
                "layout_plan": {"layout_type": "reading_task_layout", "graphic_type": "reading_structure"},
            }
        ],
    )
    issue_codes_2 = {
        detail["code"]
        for slide_report in report["slides"]
        for detail in slide_report.get("issue_details", [])
    }
    assert "graphic_data_empty" in issue_codes_2

    # 8 old task no graphic_type still export
    old_slide = {
        "slide_index": 1,
        "slide_type": "cover",
        "title": "Old Task Compatible",
        "visible_content": ["Unit 3", "Reading Lesson"],
        "teacher_notes": "notes",
        "teaching_purpose": "open class",
        "estimated_time": "2 min",
        "interaction_type": "Teacher-student interaction",
        "layout_plan": {"layout_type": "cover_layout"},
    }
    path_old = export_pptx(_task("Reading"), [old_slide])
    assert path_old.exists()
    assert has_active_graphic(old_slide) is False

    print("ROUND_019_OK")


if __name__ == "__main__":
    run()
