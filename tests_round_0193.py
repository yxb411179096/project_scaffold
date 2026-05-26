from pptx import Presentation

import services.ppt_render_service as ppt_render_service
from services.agents.layout_planner_agent import plan_layout_for_slide
from services.ppt_graphics_service import GRAPHIC_RENDERERS
from services.ppt_render_service import export_pptx


def _task(lesson_type="Reading"):
    return {
        "id": 19031,
        "course_title": "[TEST] Round193 Compact",
        "topic": "The Internet",
        "grade": "高一",
        "textbook": "人教版",
        "volume": "必修二",
        "unit": "Unit 3",
        "lesson_type": lesson_type,
        "duration": 45,
        "student_level": "中等",
        "style": "常规课",
    }


def _shape_texts(slide):
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            texts.append(shape.text_frame.text or "")
    return "\n".join(texts)


def run():
    # 1) discussion_grid should not render Roles in PPT主体
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style = {
        "primary": (21, 76, 121),
        "secondary": (51, 115, 167),
        "accent": (244, 178, 102),
        "text": (32, 43, 57),
        "muted": (102, 112, 133),
        "surface": (255, 255, 255),
        "soft": (240, 246, 252),
    }
    GRAPHIC_RENDERERS["discussion_grid"](
        slide,
        prs,
        {
            "topic": "How does the internet change our life?",
            "useful_expressions": ["I think...", "In my view...", "The text shows..."],
            "group_roles": ["Speaker", "Recorder", "Timekeeper"],
            "share_task": "Share one conclusion with evidence.",
        },
        area=(0.95, 1.85, 11.2, 4.25),
        style=style,
    )
    joined = _shape_texts(slide).lower()
    assert "roles" not in joined

    # 2) vocabulary cards should not produce Item placeholders
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    GRAPHIC_RENDERERS["vocabulary_cards"](
        slide2,
        prs,
        {
            "cards": [
                {"word": "online community", "meaning": "a group connected online", "example": "She joined an online community."},
                {"word": "social network", "meaning": "a platform for communication", "example": "The social network spreads ideas fast."},
                {"word": "digital life", "meaning": "life with internet technology", "example": "Digital life changes daily habits."},
            ]
        },
        area=(0.95, 1.85, 11.2, 4.25),
        style=style,
    )
    joined2 = _shape_texts(slide2).lower()
    assert "item 1" not in joined2 and "item 2" not in joined2

    # 3) summary graphic data should avoid hard truncation ellipsis
    summary_plan = plan_layout_for_slide(
        {"slide_type": "summary", "title": "Summary", "visible_content": ["a very very long line that should not be used directly"]},
        _task("Reading"),
    )
    branches = (summary_plan.get("graphic_data") or {}).get("branches") or []
    assert branches == ["Topic", "Skill", "Language"]
    assert not any("..." in str(x) for x in branches)

    # 4) graphic_too_dense should trigger safe compact layout (legacy body renderer not called)
    orig_renderer = ppt_render_service.LAYOUT_RENDERERS["reading_task_layout"]
    called = {"legacy": 0}

    def _counting_renderer(slide_obj, slide_data, task, prs_obj, style_obj):
        called["legacy"] += 1
        return orig_renderer(slide_obj, slide_data, task, prs_obj, style_obj)

    ppt_render_service.LAYOUT_RENDERERS["reading_task_layout"] = _counting_renderer
    try:
        dense_slide = {
            "slide_index": 1,
            "slide_type": "language_points",
            "title": "Language Points",
            "visible_content": [f"long content {i} " * 8 for i in range(6)],
            "teacher_notes": "Detailed notes " * 20,
            "teaching_purpose": "language focus",
            "estimated_time": "6 min",
            "interaction_type": "Teacher-student interaction",
            "graphic_type": "grammar_rule_chart",
            "graphic_data": {
                "examples": ["x" * 120, "y" * 120],
                "rule": "z" * 160,
                "practice": "p" * 120,
                "common_mistakes": ["m" * 120, "n" * 120],
            },
            "layout_plan": {"layout_type": "reading_task_layout", "graphic_type": "grammar_rule_chart"},
        }
        path = export_pptx(_task("Reading"), [dense_slide])
        assert path.exists()
        assert called["legacy"] == 0
    finally:
        ppt_render_service.LAYOUT_RENDERERS["reading_task_layout"] = orig_renderer

    # 5) font size should not go below 11 in graphics output
    min_size = 100
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.font.size is not None:
                    min_size = min(min_size, int(round(p.font.size.pt)))
    assert min_size >= 11

    # 6) Reading默认图形分配仅 careful_reading/vocabulary_focus/summary
    assert plan_layout_for_slide({"slide_type": "careful_reading", "visible_content": ["a"]}, _task("Reading")).get("graphic_type") == "reading_structure"
    assert plan_layout_for_slide({"slide_type": "vocabulary_focus", "visible_content": ["a"]}, _task("Reading")).get("graphic_type") == "vocabulary_cards"
    assert plan_layout_for_slide({"slide_type": "summary", "visible_content": ["a"]}, _task("Reading")).get("graphic_type") == "mindmap"
    assert plan_layout_for_slide({"slide_type": "lead_in", "visible_content": ["a"]}, _task("Reading")).get("graphic_type") == "none"
    assert plan_layout_for_slide({"slide_type": "fast_reading", "visible_content": ["a"]}, _task("Reading")).get("graphic_type") == "none"

    # 7) old task still export
    old = {
        "slide_index": 1,
        "slide_type": "cover",
        "title": "Legacy",
        "visible_content": ["Unit 3", "Reading"],
        "teacher_notes": "notes",
        "teaching_purpose": "open",
        "estimated_time": "2 min",
        "interaction_type": "Teacher-student interaction",
        "layout_plan": {"layout_type": "cover_layout"},
    }
    old_path = export_pptx(_task("Reading"), [old])
    assert old_path.exists()

    print("ROUND_0193_OK")


if __name__ == "__main__":
    run()
